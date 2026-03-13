#!/usr/bin/env python3
"""
Full integration test for the Document Conversion Service.

This script:
  1. Generates sample files for every supported document type.
  2. Tests each converter individually and shows extracted text.
  3. Tests the full processor pipeline (convert → write to S3-style output).
  4. Tests the FastAPI endpoints via TestClient.
  5. Prints a colour-coded results table.

Run from the project root:
    cd docconv-service
    python -m tests.test_full_flow
"""

import io
import json
import os
import sys
import tempfile
import textwrap
import time
import traceback
from pathlib import Path

# ── Make sure project root is on sys.path ────────────────────────────────────
PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

os.environ.setdefault("TMP_DIR", "/tmp/docconv")
os.environ.setdefault("S3_ENDPOINT_URL", "http://localhost:4566")
os.environ.setdefault("S3_OUTPUT_BUCKET", "docconv-output")
os.environ.setdefault("ENABLE_API", "true")
os.environ.setdefault("ENABLE_SQS", "false")
os.environ.setdefault("ENABLE_RABBITMQ", "false")
os.environ.setdefault("ENABLE_KAFKA", "false")
os.environ.setdefault("ENABLE_TEMPORAL", "false")
os.environ.setdefault("USE_TEMPORAL_WORKFLOWS", "false")
os.environ.setdefault("LOG_LEVEL", "WARNING")

os.makedirs("/tmp/docconv", exist_ok=True)

# ── Colour helpers ───────────────────────────────────────────────────────────
GREEN = "\033[92m"
RED = "\033[91m"
YELLOW = "\033[93m"
CYAN = "\033[96m"
BOLD = "\033[1m"
RESET = "\033[0m"

PASS = f"{GREEN}✓ PASS{RESET}"
FAIL = f"{RED}✗ FAIL{RESET}"
SKIP = f"{YELLOW}⊘ SKIP{RESET}"

# ═════════════════════════════════════════════════════════════════════════════
# SECTION 1 – SAMPLE FILE GENERATORS
# ═════════════════════════════════════════════════════════════════════════════

SAMPLE_DIR = Path("/tmp/docconv/samples")
SAMPLE_DIR.mkdir(parents=True, exist_ok=True)

SAMPLE_TEXT = (
    "The Document Conversion Service converts multiple file formats into "
    "plain text. It supports PDF, DOCX, XLSX, PPTX, HTML, RTF, ODT, TXT, "
    "CSV, and various image formats. Extracted text is uploaded to S3."
)


def make_txt() -> Path:
    p = SAMPLE_DIR / "sample.txt"
    p.write_text(SAMPLE_TEXT, encoding="utf-8")
    return p


def make_csv() -> Path:
    p = SAMPLE_DIR / "sample.csv"
    p.write_text(
        "Name,Role,Department\n"
        "Alice,Engineer,Backend\n"
        "Bob,Designer,Frontend\n"
        "Charlie,Manager,Operations\n",
        encoding="utf-8",
    )
    return p


def make_html() -> Path:
    p = SAMPLE_DIR / "sample.html"
    p.write_text(
        "<!DOCTYPE html><html><head><title>Test</title></head>"
        f"<body><h1>Document Test</h1><p>{SAMPLE_TEXT}</p>"
        "<script>var x=1;</script></body></html>",
        encoding="utf-8",
    )
    return p


def make_rtf() -> Path:
    p = SAMPLE_DIR / "sample.rtf"
    # Minimal RTF wrapper
    p.write_text(
        r"{\rtf1\ansi " + SAMPLE_TEXT + r"}",
        encoding="utf-8",
    )
    return p


def make_docx() -> Path:
    from docx import Document

    p = SAMPLE_DIR / "sample.docx"
    doc = Document()
    doc.add_heading("Conversion Test", level=1)
    doc.add_paragraph(SAMPLE_TEXT)
    # Add a table
    table = doc.add_table(rows=3, cols=2)
    table.cell(0, 0).text = "Format"
    table.cell(0, 1).text = "Status"
    table.cell(1, 0).text = "DOCX"
    table.cell(1, 1).text = "Supported"
    table.cell(2, 0).text = "PDF"
    table.cell(2, 1).text = "Supported"
    doc.save(str(p))
    return p


def make_xlsx() -> Path:
    from openpyxl import Workbook

    p = SAMPLE_DIR / "sample.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "TestSheet"
    ws.append(["Product", "Q1 Sales", "Q2 Sales"])
    ws.append(["Widget A", 1500, 2300])
    ws.append(["Widget B", 800, 1200])
    ws.append(["Widget C", 3200, 2900])
    wb.save(str(p))
    return p


def make_pptx() -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    p = SAMPLE_DIR / "sample.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Conversion Service Demo"
    slide.placeholders[1].text = SAMPLE_TEXT
    # Add notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Speaker note: demo for testing."
    # Second slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Slide Two"
    slide2.placeholders[1].text = "Additional content on slide two."
    prs.save(str(p))
    return p


def make_odt() -> Path:
    from odf.opendocument import OpenDocumentText
    from odf.text import P

    p = SAMPLE_DIR / "sample.odt"
    doc = OpenDocumentText()
    para = P(text=SAMPLE_TEXT)
    doc.text.addElement(para)
    para2 = P(text="Second paragraph in the ODT file.")
    doc.text.addElement(para2)
    doc.save(str(p))
    return p


def make_pdf() -> Path:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas as rl_canvas

    p = SAMPLE_DIR / "sample.pdf"
    c = rl_canvas.Canvas(str(p), pagesize=letter)

    # Page 1 – selectable text
    c.setFont("Helvetica", 14)
    c.drawString(72, 700, "Document Conversion Service – Test PDF")
    c.setFont("Helvetica", 11)
    # Wrap text across lines
    y = 660
    for line in textwrap.wrap(SAMPLE_TEXT, width=80):
        c.drawString(72, y, line)
        y -= 16
    c.drawString(72, y - 20, "Page 1 of 2")
    c.showPage()

    # Page 2 – more text
    c.setFont("Helvetica", 11)
    c.drawString(72, 700, "This is page two of the test PDF.")
    c.drawString(72, 680, "It contains additional searchable text content.")
    c.drawString(72, 660, "Page 2 of 2")
    c.showPage()

    c.save()
    return p


def make_image() -> Path:
    from PIL import Image, ImageDraw, ImageFont

    p = SAMPLE_DIR / "sample.png"
    img = Image.new("RGB", (600, 200), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)
    # Use default font
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 24)
    except (IOError, OSError):
        font = ImageFont.load_default()
    draw.text((30, 30), "Document Conversion", fill=(0, 0, 0), font=font)
    draw.text((30, 80), "OCR Test Image 12345", fill=(0, 0, 0), font=font)
    draw.text((30, 130), "Hello World!", fill=(0, 0, 0), font=font)
    img.save(str(p))
    return p


# Map of generators
SAMPLE_GENERATORS = {
    "txt": make_txt,
    "csv": make_csv,
    "html": make_html,
    "rtf": make_rtf,
    "docx": make_docx,
    "xlsx": make_xlsx,
    "pptx": make_pptx,
    "odt": make_odt,
    "pdf": make_pdf,
    "image": make_image,
}


# ═════════════════════════════════════════════════════════════════════════════
# SECTION 2 – CONVERTER UNIT TESTS
# ═════════════════════════════════════════════════════════════════════════════

from app.converters import (
    pdf_converter,
    docx_converter,
    xlsx_converter,
    pptx_converter,
    html_converter,
    rtf_converter,
    odt_converter,
    text_converter,
    image_converter,
)
from app.converters.dispatch import convert_document
from app.models import DocumentType, LocationType, ConversionJob

CONVERTER_MAP = {
    "txt": (text_converter, DocumentType.TXT),
    "csv": (xlsx_converter, DocumentType.CSV),
    "html": (html_converter, DocumentType.HTML),
    "rtf": (rtf_converter, DocumentType.RTF),
    "docx": (docx_converter, DocumentType.DOCX),
    "xlsx": (xlsx_converter, DocumentType.XLSX),
    "pptx": (pptx_converter, DocumentType.PPTX),
    "odt": (odt_converter, DocumentType.ODT),
    "pdf": (pdf_converter, DocumentType.PDF),
    "image": (image_converter, DocumentType.IMAGE),
}

# ═════════════════════════════════════════════════════════════════════════════
# SECTION 3 – PIPELINE TEST (processor without real S3)
# ═════════════════════════════════════════════════════════════════════════════

def test_pipeline_local(file_path: Path, doc_type: DocumentType) -> dict:
    """
    Test the convert stage of the pipeline and write output to a local file
    (no S3 needed). Returns dict with success, output_path, chars, preview.
    """
    output_path = Path("/tmp/docconv/output") / f"{file_path.stem}.txt"
    output_path.parent.mkdir(parents=True, exist_ok=True)

    total_chars = 0
    with open(output_path, "w", encoding="utf-8") as out:
        for chunk in convert_document(str(file_path), doc_type):
            out.write(chunk)
            total_chars += len(chunk)

    preview = output_path.read_text(encoding="utf-8")[:500]
    return {
        "success": True,
        "output_path": str(output_path),
        "chars": total_chars,
        "preview": preview,
    }


# ═════════════════════════════════════════════════════════════════════════════
# SECTION 4 – API TESTS
# ═════════════════════════════════════════════════════════════════════════════

def test_api():
    """
    Test the FastAPI endpoints using Starlette TestClient (no real server needed).
    We monkeypatch the S3 upload to write locally instead.
    """
    results = {}
    from fastapi.testclient import TestClient
    from app.api import app
    import app.storage as storage_mod
    import app.processor as processor_mod

    # ── Monkeypatch S3 upload to local file ──────────────────────────────
    _original_upload = storage_mod.upload_text_chunks

    def _mock_upload(text_chunks, bucket=None, key=None, job_id=None):
        out_dir = Path("/tmp/docconv/api_output")
        out_dir.mkdir(parents=True, exist_ok=True)
        import uuid
        _id = job_id or uuid.uuid4().hex[:12]
        _key = key or f"converted/{_id}.txt"
        _bucket = bucket or "docconv-output"
        out_path = out_dir / f"{_id}.txt"
        total = 0
        with open(out_path, "w") as f:
            for chunk in text_chunks:
                f.write(chunk)
                total += len(chunk)
        return _bucket, _key, total

    # Patch in BOTH modules so all call paths hit the mock
    storage_mod.upload_text_chunks = _mock_upload
    processor_mod.upload_text_chunks = _mock_upload

    client = TestClient(app)

    # ── Test /health ─────────────────────────────────────────────────────
    try:
        resp = client.get("/health")
        assert resp.status_code == 200
        data = resp.json()
        assert data["status"] == "ok"
        results["GET /health"] = {"status": PASS, "response": data}
    except Exception as e:
        results["GET /health"] = {"status": FAIL, "error": str(e)}

    # ── Test /convert/upload with TXT ────────────────────────────────────
    try:
        content = b"API upload test content.\nLine two.\nLine three."
        resp = client.post(
            "/convert/upload",
            files={"file": ("test.txt", content, "text/plain")},
            data={"document_type": "txt"},
        )
        assert resp.status_code == 200
        data = resp.json()
        assert data["success"] is True
        assert data["characters_extracted"] > 0
        results["POST /convert/upload (txt)"] = {"status": PASS, "response": data}
    except Exception as e:
        results["POST /convert/upload (txt)"] = {"status": FAIL, "error": str(e)}

    # ── Test /convert/upload with DOCX ───────────────────────────────────
    try:
        docx_path = make_docx()
        with open(docx_path, "rb") as f:
            resp = client.post(
                "/convert/upload",
                files={"file": ("sample.docx", f, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")},
                data={"document_type": "docx"},
            )
        assert resp.status_code == 200
        data = resp.json()
        assert data["success"] is True
        results["POST /convert/upload (docx)"] = {"status": PASS, "response": data}
    except Exception as e:
        results["POST /convert/upload (docx)"] = {"status": FAIL, "error": str(e)}

    # ── Test /convert/upload with PDF ────────────────────────────────────
    try:
        pdf_path = make_pdf()
        with open(pdf_path, "rb") as f:
            resp = client.post(
                "/convert/upload",
                files={"file": ("sample.pdf", f, "application/pdf")},
                data={"document_type": "pdf"},
            )
        assert resp.status_code == 200
        data = resp.json()
        assert data["success"] is True
        results["POST /convert/upload (pdf)"] = {"status": PASS, "response": data}
    except Exception as e:
        results["POST /convert/upload (pdf)"] = {"status": FAIL, "error": str(e)}

    # ── Test /convert/upload with IMAGE ──────────────────────────────────
    try:
        img_path = make_image()
        with open(img_path, "rb") as f:
            resp = client.post(
                "/convert/upload",
                files={"file": ("sample.png", f, "image/png")},
                data={"document_type": "image"},
            )
        assert resp.status_code == 200
        data = resp.json()
        assert data["success"] is True
        results["POST /convert/upload (image)"] = {"status": PASS, "response": data}
    except Exception as e:
        results["POST /convert/upload (image)"] = {"status": FAIL, "error": str(e)}

    # ── Test /convert/upload with XLSX ───────────────────────────────────
    try:
        xlsx_path = make_xlsx()
        with open(xlsx_path, "rb") as f:
            resp = client.post(
                "/convert/upload",
                files={"file": ("sample.xlsx", f, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")},
                data={"document_type": "xlsx"},
            )
        assert resp.status_code == 200
        data = resp.json()
        assert data["success"] is True
        results["POST /convert/upload (xlsx)"] = {"status": PASS, "response": data}
    except Exception as e:
        results["POST /convert/upload (xlsx)"] = {"status": FAIL, "error": str(e)}

    # ── Test validation: LOCAL location on /convert/job should fail ──────
    try:
        resp = client.post(
            "/convert/job",
            json={
                "document_type": "txt",
                "location_type": "local",
            },
        )
        assert resp.status_code == 400
        results["POST /convert/job (reject local)"] = {"status": PASS, "response": resp.json()}
    except Exception as e:
        results["POST /convert/job (reject local)"] = {"status": FAIL, "error": str(e)}

    # Restore original
    storage_mod.upload_text_chunks = _original_upload
    processor_mod.upload_text_chunks = _original_upload
    return results


# ═════════════════════════════════════════════════════════════════════════════
# SECTION 5 – MODELS VALIDATION TESTS
# ═════════════════════════════════════════════════════════════════════════════

def test_models():
    results = {}

    # Valid S3 job
    try:
        job = ConversionJob(
            job_id="test-1",
            document_type="pdf",
            location_type="s3",
            s3_bucket="my-bucket",
            s3_key="docs/file.pdf",
        )
        assert job.document_type == DocumentType.PDF
        assert job.location_type == LocationType.S3
        results["Model: valid S3 job"] = {"status": PASS}
    except Exception as e:
        results["Model: valid S3 job"] = {"status": FAIL, "error": str(e)}

    # Valid URL job with auth
    try:
        job = ConversionJob(
            document_type="html",
            location_type="url",
            url="https://example.com/doc.html",
            auth_type="bearer",
            auth_token="tok_abc123",
        )
        assert job.auth_type.value == "bearer"
        results["Model: URL job with bearer auth"] = {"status": PASS}
    except Exception as e:
        results["Model: URL job with bearer auth"] = {"status": FAIL, "error": str(e)}

    # Valid FTP job
    try:
        job = ConversionJob(
            document_type="txt",
            location_type="ftp",
            ftp_host="ftp.example.com",
            ftp_path="/data/file.txt",
            ftp_user="user",
            ftp_pass="pass",
        )
        assert job.ftp_port == 21  # default
        results["Model: FTP job defaults"] = {"status": PASS}
    except Exception as e:
        results["Model: FTP job defaults"] = {"status": FAIL, "error": str(e)}

    # Invalid document type
    try:
        from pydantic import ValidationError
        try:
            ConversionJob(document_type="mp3", location_type="url", url="http://x")
            results["Model: reject invalid doc type"] = {"status": FAIL, "error": "Should have raised"}
        except ValidationError:
            results["Model: reject invalid doc type"] = {"status": PASS}
    except Exception as e:
        results["Model: reject invalid doc type"] = {"status": FAIL, "error": str(e)}

    return results


# ═════════════════════════════════════════════════════════════════════════════
# SECTION 5b – TEMPORAL WORKFLOW TESTS (offline – no server needed)
# ═════════════════════════════════════════════════════════════════════════════

def test_temporal():
    """Test Temporal dataclasses, workflow input conversion, and activity logic."""
    results = {}

    # ── Dataclass serialization ──────────────────────────────────────────
    try:
        from app.workflows.dataclasses import (
            ConversionWorkflowInput, ConversionWorkflowOutput,
            FetchInput, FetchOutput, ConvertInput, ConvertOutput,
            UploadInput, UploadOutput, CleanupInput,
        )
        inp = ConversionWorkflowInput(
            job_id="test-temporal-1",
            document_type="pdf",
            location_type="s3",
            s3_bucket="test-bucket",
            s3_key="test.pdf",
        )
        assert inp.job_id == "test-temporal-1"
        assert inp.document_type == "pdf"
        assert inp.ftp_port == 21  # default
        assert inp.auth_type == "none"
        results["Temporal: dataclass creation"] = {"status": PASS}
    except Exception as e:
        results["Temporal: dataclass creation"] = {"status": FAIL, "error": str(e)}

    # ── Workflow input from ConversionJob ─────────────────────────────────
    try:
        from app.workflows.client import _job_to_workflow_input
        from app.models import ConversionJob
        job = ConversionJob(
            job_id="wf-test-1",
            document_type="docx",
            location_type="url",
            url="https://example.com/doc.docx",
            auth_type="bearer",
            auth_token="tok_abc",
        )
        wf_inp = _job_to_workflow_input(job)
        assert wf_inp.job_id == "wf-test-1"
        assert wf_inp.document_type == "docx"
        assert wf_inp.location_type == "url"
        assert wf_inp.url == "https://example.com/doc.docx"
        assert wf_inp.auth_type == "bearer"
        assert wf_inp.auth_token == "tok_abc"
        results["Temporal: job → workflow input"] = {"status": PASS}
    except Exception as e:
        results["Temporal: job → workflow input"] = {"status": FAIL, "error": str(e)}

    # ── Workflow output → ConversionResult ────────────────────────────────
    try:
        from app.workflows.client import _workflow_output_to_result
        from app.workflows.dataclasses import ConversionWorkflowOutput
        wf_out = ConversionWorkflowOutput(
            job_id="wf-test-1",
            success=True,
            output_bucket="docconv-output",
            output_key="converted/wf-test-1.txt",
            total_chars=1234,
            pages_processed=5,
            images_extracted=2,
        )
        result = _workflow_output_to_result(wf_out)
        assert result.success is True
        assert result.output_bucket == "docconv-output"
        assert result.characters_extracted == 1234
        assert result.job_id == "wf-test-1"
        results["Temporal: workflow output → result"] = {"status": PASS}
    except Exception as e:
        results["Temporal: workflow output → result"] = {"status": FAIL, "error": str(e)}

    # ── Child workflow registry ───────────────────────────────────────────
    try:
        from app.workflows.document_workflows import CHILD_WORKFLOW_MAP, ALL_CHILD_WORKFLOWS
        # All 10 document types should be mapped
        expected = {"pdf", "docx", "xlsx", "csv", "pptx", "html", "rtf", "odt", "txt", "image"}
        assert set(CHILD_WORKFLOW_MAP.keys()) == expected, \
            f"Missing: {expected - set(CHILD_WORKFLOW_MAP.keys())}"
        assert len(ALL_CHILD_WORKFLOWS) == 9  # csv shares xlsx workflow
        results["Temporal: child workflow registry"] = {"status": PASS}
    except Exception as e:
        results["Temporal: child workflow registry"] = {"status": FAIL, "error": str(e)}

    # ── Activity registry ────────────────────────────────────────────────
    try:
        from app.workflows.activities import ALL_ACTIVITIES
        assert len(ALL_ACTIVITIES) == 12  # fetch + 9 convert + upload + cleanup
        # Extended activities
        from app.workflows.activities_ext import ALL_EXTENDED_ACTIVITIES
        assert len(ALL_EXTENDED_ACTIVITIES) == 10  # scan, move, validate, analyze, metadata, webhook, ocr, split, cleanup, health
        # Verify all have temporal activity definitions
        for act in ALL_ACTIVITIES:
            assert hasattr(act, "__temporal_activity_definition"), \
                f"{act.__name__} is not a Temporal activity"
        results["Temporal: activity registry"] = {"status": PASS}
    except Exception as e:
        results["Temporal: activity registry"] = {"status": FAIL, "error": str(e)}

    # ── All workflow classes have definitions ─────────────────────────────
    try:
        from app.workflows.worker import ALL_WORKFLOWS, COMBINED_ACTIVITIES
        assert len(ALL_WORKFLOWS) == 17  # 2 core + 9 child + 6 extended
        assert len(COMBINED_ACTIVITIES) == 22  # 12 core + 10 extended
        for wf_cls in ALL_WORKFLOWS:
            assert hasattr(wf_cls, "__temporal_workflow_definition"), \
                f"{wf_cls.__name__} is not a Temporal workflow"
        results["Temporal: workflow class registry"] = {"status": PASS}
    except Exception as e:
        results["Temporal: workflow class registry"] = {"status": FAIL, "error": str(e)}

    # ── Extended dataclasses ─────────────────────────────────────────────
    try:
        from app.workflows.dataclasses_ext import (
            S3FolderWatchInput, PipelineWorkflowInput,
            MultiFormatInput, RetryEscalationInput,
            WebhookNotificationWorkflowInput,
            ScheduledCleanupInput, ValidateInput, AnalyzeTextInput,
        )
        # Verify defaults
        s3w = S3FolderWatchInput(bucket="test")
        assert s3w.prefix == "inbox/"
        assert s3w.max_files == 100
        pipe = PipelineWorkflowInput(job_id="x", document_type="pdf", location_type="s3")
        assert pipe.enable_validation is True
        assert pipe.enable_analytics is True
        mf = MultiFormatInput(job_id="x", document_type="pdf", location_type="url")
        assert mf.produce_pages is True
        re_inp = RetryEscalationInput(job_id="x", document_type="pdf", location_type="s3")
        assert re_inp.min_chars_threshold == 50
        assert re_inp.escalation_dpi == 400
        wh = WebhookNotificationWorkflowInput(job_id="x", document_type="pdf", location_type="url")
        assert wh.on_start_webhook is None
        sc = ScheduledCleanupInput()
        assert sc.max_age_hours == 24
        results["Temporal: extended dataclasses"] = {"status": PASS}
    except Exception as e:
        results["Temporal: extended dataclasses"] = {"status": FAIL, "error": str(e)}

    # ── Extended workflow imports ─────────────────────────────────────────
    try:
        from app.workflows.pipeline_workflow import DocumentPipelineWorkflow
        from app.workflows.s3_watch_workflow import S3FolderWatchWorkflow
        from app.workflows.webhook_workflow import WebhookNotificationWorkflow
        from app.workflows.multi_output_workflow import MultiFormatOutputWorkflow
        from app.workflows.retry_escalation_workflow import RetryEscalationWorkflow
        from app.workflows.scheduled_workflow import ScheduledMaintenanceWorkflow
        # Verify all have temporal definitions
        for wf in [DocumentPipelineWorkflow, S3FolderWatchWorkflow,
                    WebhookNotificationWorkflow, MultiFormatOutputWorkflow,
                    RetryEscalationWorkflow, ScheduledMaintenanceWorkflow]:
            assert hasattr(wf, "__temporal_workflow_definition"), f"{wf.__name__} missing defn"
        results["Temporal: extended workflow imports"] = {"status": PASS}
    except Exception as e:
        results["Temporal: extended workflow imports"] = {"status": FAIL, "error": str(e)}

    # ── Extended activity imports ─────────────────────────────────────────
    try:
        from app.workflows.activities_ext import (
            scan_s3_prefix_activity, validate_document_activity,
            analyze_text_activity, send_webhook_activity,
            enhanced_ocr_convert_activity, split_text_by_pages_activity,
            scheduled_tmp_cleanup_activity, check_s3_health_activity,
            generate_metadata_json_activity, move_s3_object_activity,
        )
        for act in [scan_s3_prefix_activity, validate_document_activity,
                     analyze_text_activity, send_webhook_activity,
                     enhanced_ocr_convert_activity, split_text_by_pages_activity,
                     scheduled_tmp_cleanup_activity, check_s3_health_activity,
                     generate_metadata_json_activity, move_s3_object_activity]:
            assert hasattr(act, "__temporal_activity_definition"), f"{act.__name__} missing defn"
        results["Temporal: extended activity imports"] = {"status": PASS}
    except Exception as e:
        results["Temporal: extended activity imports"] = {"status": FAIL, "error": str(e)}

    # ── Direct fallback when Temporal disabled ───────────────────────────
    try:
        from config.settings import settings
        assert settings.enable_temporal is False  # disabled in test env
        assert settings.use_temporal_workflows is False
        # process_job should use direct path
        from app.models import ConversionJob
        from app.processor import process_job
        import app.storage as storage_mod
        import app.processor as processor_mod

        _orig = storage_mod.upload_text_chunks
        def _mock(text_chunks, bucket=None, key=None, job_id=None):
            total = 0
            for chunk in text_chunks:
                total += len(chunk)
            return bucket or "b", key or "k", total
        storage_mod.upload_text_chunks = _mock
        processor_mod.upload_text_chunks = _mock

        job = ConversionJob(
            job_id="fallback-test",
            document_type="txt",
            location_type="local",
        )
        # Create a fresh temp file for this test
        import shutil
        tmp = os.path.join("/tmp/docconv", "fallback_test.txt")
        with open(tmp, "w") as f:
            f.write("Fallback test content.\nLine two.")
        r = process_job(job, local_file_path=tmp)
        assert r.success is True
        assert r.characters_extracted > 0
        results["Temporal: direct fallback works"] = {"status": PASS}

        storage_mod.upload_text_chunks = _orig
        processor_mod.upload_text_chunks = _orig
    except Exception as e:
        results["Temporal: direct fallback works"] = {"status": FAIL, "error": str(e)}

    return results


# ═════════════════════════════════════════════════════════════════════════════
# SECTION 6 – RUNNER
# ═════════════════════════════════════════════════════════════════════════════

def print_header(title: str):
    print(f"\n{BOLD}{CYAN}{'═' * 70}")
    print(f"  {title}")
    print(f"{'═' * 70}{RESET}\n")


def print_divider():
    print(f"{CYAN}{'─' * 70}{RESET}")


def truncate(text: str, max_len: int = 300) -> str:
    text = text.replace("\n", " ↵ ")
    if len(text) > max_len:
        return text[:max_len] + "…"
    return text


def main():
    total_pass = 0
    total_fail = 0
    total_skip = 0
    all_results = []  # (section, name, status, detail)

    # ── 1. Generate sample files ─────────────────────────────────────────
    print_header("1. GENERATING SAMPLE FILES")
    sample_paths = {}
    for fmt, gen_fn in SAMPLE_GENERATORS.items():
        try:
            path = gen_fn()
            size = path.stat().st_size
            sample_paths[fmt] = path
            print(f"  {GREEN}✓{RESET} {fmt:8s}  →  {path.name}  ({size:,} bytes)")
        except Exception as e:
            print(f"  {RED}✗{RESET} {fmt:8s}  →  FAILED: {e}")

    # ── 2. Test each converter ───────────────────────────────────────────
    print_header("2. CONVERTER UNIT TESTS")
    for fmt, (module, doc_type) in CONVERTER_MAP.items():
        path = sample_paths.get(fmt)
        if not path:
            print(f"  {SKIP}  {fmt:8s}  (no sample file)")
            total_skip += 1
            all_results.append(("Converter", fmt, "SKIP", "no sample"))
            continue
        try:
            chunks = list(module.convert_to_text(str(path)))
            full_text = "".join(chunks)
            n_chunks = len(chunks)
            n_chars = len(full_text)
            assert n_chars > 0, "No text extracted"
            print(f"  {PASS}  {fmt:8s}  chunks={n_chunks:3d}  chars={n_chars:5,d}  preview: {truncate(full_text, 120)}")
            total_pass += 1
            all_results.append(("Converter", fmt, "PASS", f"{n_chars} chars"))
        except Exception as e:
            print(f"  {FAIL}  {fmt:8s}  →  {e}")
            traceback.print_exc()
            total_fail += 1
            all_results.append(("Converter", fmt, "FAIL", str(e)))

    # ── 3. Test converter dispatcher ─────────────────────────────────────
    print_header("3. CONVERTER DISPATCHER TESTS")
    for fmt, (_, doc_type) in CONVERTER_MAP.items():
        path = sample_paths.get(fmt)
        if not path:
            continue
        try:
            chunks = list(convert_document(str(path), doc_type))
            full_text = "".join(chunks)
            assert len(full_text) > 0
            print(f"  {PASS}  dispatch({doc_type.value:8s})  chars={len(full_text):5,d}")
            total_pass += 1
            all_results.append(("Dispatch", doc_type.value, "PASS", f"{len(full_text)} chars"))
        except Exception as e:
            print(f"  {FAIL}  dispatch({doc_type.value:8s})  →  {e}")
            total_fail += 1
            all_results.append(("Dispatch", doc_type.value, "FAIL", str(e)))

    # ── 4. Test full pipeline (local, no S3) ─────────────────────────────
    print_header("4. FULL PIPELINE TESTS (local output)")
    for fmt, (_, doc_type) in CONVERTER_MAP.items():
        path = sample_paths.get(fmt)
        if not path:
            continue
        try:
            result = test_pipeline_local(path, doc_type)
            print(f"  {PASS}  pipeline({doc_type.value:8s})  chars={result['chars']:5,d}  out={result['output_path']}")
            total_pass += 1
            all_results.append(("Pipeline", doc_type.value, "PASS", f"{result['chars']} chars"))
        except Exception as e:
            print(f"  {FAIL}  pipeline({doc_type.value:8s})  →  {e}")
            total_fail += 1
            all_results.append(("Pipeline", doc_type.value, "FAIL", str(e)))

    # ── 5. Model validation tests ────────────────────────────────────────
    print_header("5. MODEL VALIDATION TESTS")
    model_results = test_models()
    for name, info in model_results.items():
        status = info["status"]
        if "PASS" in status:
            total_pass += 1
            all_results.append(("Models", name, "PASS", ""))
        else:
            total_fail += 1
            all_results.append(("Models", name, "FAIL", info.get("error", "")))
        detail = info.get("error", "")
        extra = f"  ({detail})" if detail else ""
        print(f"  {status}  {name}{extra}")

    # ── 6. API tests ─────────────────────────────────────────────────────
    print_header("6. REST API TESTS")
    api_results = test_api()
    for name, info in api_results.items():
        status = info["status"]
        if "PASS" in status:
            total_pass += 1
            all_results.append(("API", name, "PASS", ""))
        else:
            total_fail += 1
            all_results.append(("API", name, "FAIL", info.get("error", "")))
        detail = ""
        if "response" in info:
            resp = info["response"]
            if isinstance(resp, dict):
                # Show key fields
                chars = resp.get("characters_extracted", "")
                success = resp.get("success", "")
                if chars != "":
                    detail = f"  chars={chars}"
                elif success != "":
                    detail = f"  success={success}"
        if "error" in info:
            detail = f"  ERROR: {info['error']}"
        print(f"  {status}  {name}{detail}")

    # ── 6b. Temporal tests ───────────────────────────────────────────────
    print_header("6b. TEMPORAL WORKFLOW TESTS (offline)")
    temporal_results = test_temporal()
    for name, info in temporal_results.items():
        status = info["status"]
        if "PASS" in status:
            total_pass += 1
            all_results.append(("Temporal", name, "PASS", ""))
        elif "SKIP" in status:
            total_skip += 1
            all_results.append(("Temporal", name, "SKIP", ""))
        else:
            total_fail += 1
            all_results.append(("Temporal", name, "FAIL", info.get("error", "")))
        detail = info.get("error", "")
        extra = f"  ({detail})" if detail else ""
        print(f"  {status}  {name}{extra}")

    # ── 7. Show converted text previews ──────────────────────────────────
    print_header("7. CONVERTED TEXT PREVIEWS")
    output_dir = Path("/tmp/docconv/output")
    if output_dir.exists():
        for txt_file in sorted(output_dir.glob("*.txt")):
            content = txt_file.read_text(encoding="utf-8")
            print_divider()
            print(f"  {BOLD}{txt_file.name}{RESET}  ({len(content):,} chars)")
            print()
            # Show first 500 chars indented
            preview = content[:500]
            for line in preview.splitlines():
                print(f"    {line}")
            if len(content) > 500:
                print(f"    {YELLOW}... ({len(content) - 500:,} more chars){RESET}")
            print()

    # ── 8. Summary table ─────────────────────────────────────────────────
    print_header("8. RESULTS SUMMARY")
    print(f"  {'Section':<12s}  {'Test':<35s}  {'Status':<8s}  {'Detail'}")
    print(f"  {'─'*12}  {'─'*35}  {'─'*8}  {'─'*30}")
    for section, name, status, detail in all_results:
        if status == "PASS":
            sc = f"{GREEN}PASS{RESET}"
        elif status == "FAIL":
            sc = f"{RED}FAIL{RESET}"
        else:
            sc = f"{YELLOW}SKIP{RESET}"
        print(f"  {section:<12s}  {name:<35s}  {sc:<17s}  {detail}")

    print()
    print(f"  {BOLD}Total: {total_pass} passed, {total_fail} failed, {total_skip} skipped{RESET}")
    print()

    if total_fail == 0:
        print(f"  {GREEN}{BOLD}{'═' * 50}")
        print(f"  ★  ALL TESTS PASSED  ★")
        print(f"  {'═' * 50}{RESET}")
    else:
        print(f"  {RED}{BOLD}{'═' * 50}")
        print(f"  ✗  {total_fail} TEST(S) FAILED")
        print(f"  {'═' * 50}{RESET}")

    return 0 if total_fail == 0 else 1


if __name__ == "__main__":
    sys.exit(main())
