#!/usr/bin/env python3
"""
verify_all.py – End-to-end verification against the live Docker stack.

Tests EVERY workflow, verifies they appear in Temporal, checks S3 output.

Prerequisites:
    docker compose up --build -d
    # Wait ~60s for healthy
    pip install boto3 requests reportlab python-docx openpyxl python-pptx Pillow

Usage:
    python scripts/verify_all.py
    python scripts/verify_all.py --api-base http://myhost:8080
"""

import argparse
import io
import json
import os
import sys
import textwrap
import time
import traceback

# ─── Config ──────────────────────────────────────────────────────────────────

API = "http://localhost:8080"
S3_EP = "http://localhost:4566"
IN_BUCKET = "docconv-input"
OUT_BUCKET = "docconv-output"
AWS = dict(endpoint_url=S3_EP, aws_access_key_id="test",
           aws_secret_access_key="test", region_name="us-east-1")

# ─── Colors ──────────────────────────────────────────────────────────────────
G="\033[92m";R="\033[91m";C="\033[96m";Y="\033[93m";B="\033[1m";D="\033[2m";N="\033[0m"
PASS_COUNT=0; FAIL_COUNT=0; RESULTS=[]

def ok(msg):
    global PASS_COUNT; PASS_COUNT+=1; RESULTS.append(("PASS",msg))
    print(f"  {G}✓ PASS{N}  {msg}")

def fail(msg, err=""):
    global FAIL_COUNT; FAIL_COUNT+=1; RESULTS.append(("FAIL",msg))
    print(f"  {R}✗ FAIL{N}  {msg}")
    if err: print(f"         {D}{str(err)[:200]}{N}")

def header(t):
    print(f"\n{B}{C}{'═'*70}\n  {t}\n{'═'*70}{N}\n")

def subhead(t): print(f"\n  {B}{t}{N}")

# ─── Imports ─────────────────────────────────────────────────────────────────
try:
    import boto3
    import requests
except ImportError:
    print("pip install boto3 requests reportlab python-docx openpyxl python-pptx Pillow")
    sys.exit(1)

# ─── S3 Helpers ──────────────────────────────────────────────────────────────
def s3(): return boto3.client("s3", **AWS)

def s3_put(bucket, key, data):
    s3().put_object(Bucket=bucket, Key=key, Body=data)

def s3_get(bucket, key):
    try: return s3().get_object(Bucket=bucket, Key=key)["Body"].read().decode()
    except: return None

def s3_ls(bucket, prefix=""):
    try:
        r = s3().list_objects_v2(Bucket=bucket, Prefix=prefix)
        return [o["Key"] for o in r.get("Contents",[])]
    except: return []

def s3_ensure_bucket(name):
    try: s3().head_bucket(Bucket=name)
    except: s3().create_bucket(Bucket=name)

# ─── API Helpers ─────────────────────────────────────────────────────────────
def api_get(path):
    try:
        r = requests.get(f"{API}{path}", timeout=30)
        return r.status_code, r.json()
    except Exception as e: return 0, {"error": str(e)}

def api_post(path, json_body=None, files=None, data=None):
    try:
        if files:
            r = requests.post(f"{API}{path}", files=files, data=data, timeout=180)
        else:
            r = requests.post(f"{API}{path}", json=json_body, timeout=180)
        return r.status_code, r.json()
    except Exception as e: return 0, {"error": str(e)}

def poll_workflow(wf_id, max_wait=120):
    """Poll Temporal workflow status until done or timeout."""
    t0 = time.time()
    last = ""
    while time.time()-t0 < max_wait:
        code, data = api_get(f"/workflow/{wf_id}/status")
        if code == 200:
            st = data.get("status","")
            cs = data.get("custom_status","")
            if cs != last:
                print(f"    {D}[{time.time()-t0:.0f}s] {cs or st}{N}")
                last = cs
            if st in ("COMPLETED","FAILED","TERMINATED","CANCELLED","TIMED_OUT"):
                return st, data
        time.sleep(2)
    return "TIMEOUT", {}

def get_temporal_workflows(limit=50):
    """List workflows from Temporal via API."""
    code, data = api_get(f"/workflows/recent?limit={limit}")
    if code == 200:
        return data.get("workflows", [])
    return []

# ─── Sample File Generators ─────────────────────────────────────────────────
def make_pdf():
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=letter)
    c.setFont("Helvetica",14); c.drawString(72,700,"Test PDF – Page 1")
    c.setFont("Helvetica",11)
    for i,line in enumerate(textwrap.wrap(
        "The document conversion service extracts text from PDF files. "
        "It supports selectable text and OCR for scanned pages. "
        "Embedded images are extracted and processed separately.",80)):
        c.drawString(72, 660-i*16, line)
    c.showPage()
    c.setFont("Helvetica",11); c.drawString(72,700,"Page two content here.")
    c.drawString(72,680,"Additional text on the second page.")
    c.showPage(); c.save()
    return buf.getvalue()

def make_docx():
    from docx import Document
    d = Document(); d.add_heading("Test DOCX",1)
    d.add_paragraph("This Word document tests the DOCX converter.")
    t = d.add_table(rows=2,cols=2)
    t.cell(0,0).text="Name"; t.cell(0,1).text="Value"
    t.cell(1,0).text="Test"; t.cell(1,1).text="OK"
    buf = io.BytesIO(); d.save(buf); return buf.getvalue()

def make_xlsx():
    from openpyxl import Workbook
    wb = Workbook(); ws = wb.active; ws.title="Data"
    ws.append(["Month","Sales"]); ws.append(["Jan",5000]); ws.append(["Feb",6200])
    buf = io.BytesIO(); wb.save(buf); return buf.getvalue()

def make_pptx():
    from pptx import Presentation
    p = Presentation(); s = p.slides.add_slide(p.slide_layouts[1])
    s.shapes.title.text = "Test Slide"
    s.placeholders[1].text = "Slide body text for testing."
    s.notes_slide.notes_text_frame.text = "Speaker notes here."
    buf = io.BytesIO(); p.save(buf); return buf.getvalue()

def make_html():
    return b"<html><head><title>Test</title></head><body><h1>Test HTML</h1><p>Body paragraph for conversion testing.</p><script>var x=1;</script></body></html>"

def make_txt():
    return b"Plain text file for conversion.\nLine two.\nLine three with more content."

def make_image():
    from PIL import Image, ImageDraw, ImageFont
    img = Image.new("RGB",(500,150),(255,255,255)); draw = ImageDraw.Draw(img)
    try: font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",22)
    except: font = ImageFont.load_default()
    draw.text((20,20),"Invoice #12345",fill=(0,0,0),font=font)
    draw.text((20,60),"Amount: $9,876",fill=(0,0,0),font=font)
    draw.text((20,100),"OCR Test Line",fill=(0,0,0),font=font)
    buf = io.BytesIO(); img.save(buf,format="PNG"); return buf.getvalue()


# ═════════════════════════════════════════════════════════════════════════════
# TEST SECTIONS
# ═════════════════════════════════════════════════════════════════════════════

def test_infrastructure():
    header("1. Infrastructure Connectivity")

    # API health
    code, data = api_get("/health")
    if code == 200 and data.get("status") == "ok":
        ok(f"API healthy at {API}")
        for k in ["temporal_enabled","temporal_workflows_active","sqs_enabled","rabbitmq_enabled","kafka_enabled"]:
            v = data.get(k)
            status = f"{G}✓{N}" if v else f"{Y}○{N}"
            print(f"    {status} {k}: {v}")
    else:
        fail("API not reachable", data)
        return False

    # S3 / LocalStack
    try:
        s3_ensure_bucket(IN_BUCKET); s3_ensure_bucket(OUT_BUCKET)
        s3_put(IN_BUCKET, "_healthcheck", b"ok")
        assert s3_get(IN_BUCKET, "_healthcheck") == "ok"
        ok("S3 (LocalStack) connected – buckets ready")
    except Exception as e:
        fail("S3 (LocalStack) connection failed", e); return False

    # Temporal
    if data.get("temporal_enabled"):
        code2, data2 = api_get("/workflows/recent?limit=1")
        if code2 == 200:
            ok(f"Temporal connected via API – {data.get('temporal_host')}")
        else:
            fail("Temporal API endpoint failed", data2)
    else:
        fail("Temporal not enabled in config")

    return True


def test_upload_samples():
    header("2. Upload Sample Files to S3")
    samples = {
        "verify/report.pdf":  make_pdf(),
        "verify/doc.docx":    make_docx(),
        "verify/data.xlsx":   make_xlsx(),
        "verify/slides.pptx": make_pptx(),
        "verify/page.html":   make_html(),
        "verify/notes.txt":   make_txt(),
        "verify/scan.png":    make_image(),
        "inbox/auto1.txt":    b"Auto-convert file one.\nSecond line.",
        "inbox/auto2.html":   b"<html><body><p>Auto convert file two.</p></body></html>",
        "inbox/auto3.png":    make_image(),
    }
    for key, data in samples.items():
        s3_put(IN_BUCKET, key, data)
        ok(f"s3://{IN_BUCKET}/{key}  ({len(data):,} bytes)")
    return True


def test_direct_upload():
    """Test POST /convert/upload – file uploads that go through Temporal."""
    header("3. Direct File Upload (should appear in Temporal)")

    uploads = [
        ("test.pdf",  "pdf",   make_pdf()),
        ("test.docx", "docx",  make_docx()),
        ("test.xlsx", "xlsx",  make_xlsx()),
        ("test.pptx", "pptx",  make_pptx()),
        ("test.html", "html",  make_html()),
        ("test.txt",  "txt",   make_txt()),
        ("test.png",  "image", make_image()),
    ]

    for fname, dtype, data in uploads:
        code, result = api_post("/convert/upload",
            files={"file": (fname, data)}, data={"document_type": dtype})
        if code == 200 and result.get("success"):
            chars = result.get("characters_extracted", 0)
            okey = result.get("output_key","")
            ok(f"Upload {fname:12s}  chars={chars:5d}  → {okey}")
            # Verify text in S3
            text = s3_get(result["output_bucket"], okey)
            if text and len(text) > 0:
                preview = text[:80].replace("\n"," ↵ ")
                print(f"    {D}Preview: {preview}...{N}")
            else:
                fail(f"Upload {fname} – output empty in S3")
        else:
            fail(f"Upload {fname}", result.get("error",""))


def test_s3_conversion():
    """Test POST /convert/job with S3 sources – all go through Temporal."""
    header("4. S3 → Temporal Workflow Conversions")

    jobs = [
        ("verify-pdf",  "pdf",   "verify/report.pdf"),
        ("verify-docx", "docx",  "verify/doc.docx"),
        ("verify-xlsx", "xlsx",  "verify/data.xlsx"),
        ("verify-pptx", "pptx",  "verify/slides.pptx"),
        ("verify-html", "html",  "verify/page.html"),
        ("verify-txt",  "txt",   "verify/notes.txt"),
        ("verify-img",  "image", "verify/scan.png"),
    ]

    for job_id, dtype, s3key in jobs:
        code, result = api_post("/convert/job", {
            "job_id": job_id,
            "document_type": dtype,
            "location_type": "s3",
            "s3_bucket": IN_BUCKET,
            "s3_key": s3key,
            "s3_endpoint_url": "http://localstack:4566",
        })
        if code == 200 and result.get("success"):
            chars = result.get("characters_extracted", 0)
            ok(f"S3 {dtype:6s}  chars={chars:5d}  wf=docconv-{job_id}")
        else:
            fail(f"S3 {dtype}", result.get("error",""))


def test_pipeline_workflow():
    header("5. Pipeline Workflow (validate → convert → analyze → metadata)")

    code, result = api_post("/workflow/pipeline", {
        "job_id": "verify-pipe",
        "document_type": "pdf",
        "location_type": "s3",
        "s3_bucket": IN_BUCKET,
        "s3_key": "verify/report.pdf",
        "s3_endpoint_url": "http://localstack:4566",
        "output_prefix": "pipeline/",
        "enable_validation": True,
        "enable_analytics": True,
        "enable_metadata_output": True,
    })

    if not result.get("workflow_id"):
        fail("Pipeline start failed", result); return

    wf_id = result["workflow_id"]
    ok(f"Pipeline started: {wf_id}")
    subhead("Polling...")
    status, _ = poll_workflow(wf_id, max_wait=120)

    if status == "COMPLETED":
        ok(f"Pipeline completed: {wf_id}")
        # Check outputs
        keys = s3_ls(OUT_BUCKET, "pipeline/verify-pipe/")
        for k in keys:
            print(f"    {G}→{N} s3://{OUT_BUCKET}/{k}")
        if any(k.endswith(".txt") for k in keys):
            ok("Pipeline text output found")
        else:
            fail("Pipeline text output missing")
        if any(k.endswith(".json") for k in keys):
            ok("Pipeline metadata JSON found")
            meta = s3_get(OUT_BUCKET, [k for k in keys if k.endswith(".json")][0])
            if meta:
                md = json.loads(meta)
                print(f"    {D}Analytics: words={md.get('analytics',{}).get('total_words',0)} "
                      f"lang={md.get('analytics',{}).get('language_hint','?')}{N}")
        else:
            fail("Pipeline metadata JSON missing")
    else:
        fail(f"Pipeline ended with status: {status}")


def test_s3_watch_workflow():
    header("6. S3 Folder Watch Workflow (scan inbox/ → convert all)")

    code, result = api_post("/workflow/s3-watch", {
        "bucket": IN_BUCKET,
        "prefix": "inbox/",
        "output_prefix": "watch/",
        "s3_endpoint_url": "http://localstack:4566",
        "move_processed_to": "processed/",
    })

    if not result.get("workflow_id"):
        fail("S3 watch start failed", result); return

    wf_id = result["workflow_id"]
    ok(f"S3 watch started: {wf_id}")
    subhead("Polling...")
    status, _ = poll_workflow(wf_id, max_wait=180)

    if status == "COMPLETED":
        ok(f"S3 watch completed: {wf_id}")
        watch_keys = s3_ls(OUT_BUCKET, "watch/")
        proc_keys = s3_ls(IN_BUCKET, "processed/")
        print(f"    Outputs: {len(watch_keys)} files converted")
        print(f"    Moved:   {len(proc_keys)} files to processed/")
        if len(watch_keys) > 0: ok(f"Watch produced {len(watch_keys)} output files")
        else: fail("Watch produced no output")
    else:
        fail(f"S3 watch ended with status: {status}")


def test_multi_format_workflow():
    header("7. Multi-Format Output Workflow")

    code, result = api_post("/workflow/multi-format", {
        "job_id": "verify-multi",
        "document_type": "pdf",
        "location_type": "s3",
        "s3_bucket": IN_BUCKET,
        "s3_key": "verify/report.pdf",
        "s3_endpoint_url": "http://localstack:4566",
        "output_prefix": "multi/",
        "produce_full_text": True,
        "produce_pages": True,
        "produce_metadata_json": True,
        "produce_analytics": True,
    })

    if not result.get("workflow_id"):
        fail("Multi-format start failed", result); return

    wf_id = result["workflow_id"]
    ok(f"Multi-format started: {wf_id}")
    subhead("Polling...")
    status, _ = poll_workflow(wf_id, max_wait=120)

    if status == "COMPLETED":
        ok(f"Multi-format completed: {wf_id}")
        keys = sorted(s3_ls(OUT_BUCKET, "multi/verify-multi/"))
        for k in keys:
            print(f"    {G}→{N} s3://{OUT_BUCKET}/{k}")
        has_text = any("full_text" in k for k in keys)
        has_pages = any("page_" in k for k in keys)
        has_meta = any("metadata" in k for k in keys)
        if has_text: ok("Full text output found")
        else: fail("Full text missing")
        if has_pages: ok(f"Page splits found ({sum(1 for k in keys if 'page_' in k)} pages)")
        else: fail("Page splits missing")
        if has_meta: ok("Metadata JSON found")
        else: fail("Metadata JSON missing")
    else:
        fail(f"Multi-format ended with status: {status}")


def test_webhook_workflow():
    header("8. Webhook Notification Workflow")

    code, result = api_post("/workflow/webhook-convert", {
        "job_id": "verify-wh",
        "document_type": "html",
        "location_type": "s3",
        "s3_bucket": IN_BUCKET,
        "s3_key": "verify/page.html",
        "s3_endpoint_url": "http://localstack:4566",
        "on_start_webhook": "https://httpbin.org/post",
        "on_complete_webhook": "https://httpbin.org/post",
    })

    if not result.get("workflow_id"):
        fail("Webhook workflow start failed", result); return

    wf_id = result["workflow_id"]
    ok(f"Webhook workflow started: {wf_id}")
    subhead("Polling...")
    status, _ = poll_workflow(wf_id, max_wait=120)

    if status == "COMPLETED":
        ok(f"Webhook workflow completed: {wf_id}")
    else:
        fail(f"Webhook workflow ended: {status}")


def test_retry_escalation():
    header("9. Retry Escalation Workflow")

    code, result = api_post("/workflow/retry-escalation", {
        "job_id": "verify-esc",
        "document_type": "image",
        "location_type": "s3",
        "s3_bucket": IN_BUCKET,
        "s3_key": "verify/scan.png",
        "s3_endpoint_url": "http://localstack:4566",
        "min_chars_threshold": 10,
        "escalation_dpi": 400,
    })

    if not result.get("workflow_id"):
        fail("Escalation start failed", result); return

    wf_id = result["workflow_id"]
    ok(f"Escalation started: {wf_id}")
    subhead("Polling...")
    status, _ = poll_workflow(wf_id, max_wait=180)

    if status == "COMPLETED":
        ok(f"Escalation completed: {wf_id}")
    else:
        fail(f"Escalation ended: {status}")


def test_maintenance():
    header("10. Maintenance Workflow")

    code, result = api_post("/workflow/maintenance", {
        "tmp_dir": "/tmp/docconv",
        "max_age_hours": 24,
    })

    if not result.get("workflow_id"):
        fail("Maintenance start failed", result); return

    wf_id = result["workflow_id"]
    ok(f"Maintenance started: {wf_id}")
    subhead("Polling...")
    status, _ = poll_workflow(wf_id, max_wait=30)

    if status == "COMPLETED":
        ok(f"Maintenance completed: {wf_id}")
    else:
        fail(f"Maintenance ended: {status}")


def test_temporal_dashboard():
    """Verify that all workflows appear in the Temporal dashboard."""
    header("11. Temporal Dashboard Verification")

    subhead("Fetching workflow list from Temporal...")
    time.sleep(3)  # let workflows propagate
    workflows = get_temporal_workflows(limit=50)

    if not workflows:
        fail("No workflows found in Temporal")
        return

    print(f"  Found {len(workflows)} workflows in Temporal:\n")

    # Group by type
    by_type = {}
    for wf in workflows:
        wf_type = wf.get("workflow_type", "unknown")
        if wf_type not in by_type:
            by_type[wf_type] = []
        by_type[wf_type].append(wf)

    for wf_type, wfs in sorted(by_type.items()):
        print(f"    {B}{wf_type}{N}  ({len(wfs)} executions)")
        for wf in wfs[:5]:  # show first 5
            st = wf.get("status","?")
            color = G if st=="COMPLETED" else R if st=="FAILED" else Y
            print(f"      {color}{st:12s}{N}  {wf['id']}")
        if len(wfs) > 5:
            print(f"      {D}... and {len(wfs)-5} more{N}")
        print()

    # ── Verify specific workflow types appeared ──────────────────────────
    expected_types = {
        "DocumentConversionWorkflow":  "Basic/S3/Upload conversions",
        "DocumentPipelineWorkflow":    "Pipeline workflow",
        "S3FolderWatchWorkflow":       "S3 folder watch",
        "MultiFormatOutputWorkflow":   "Multi-format output",
        "WebhookNotificationWorkflow": "Webhook notification",
        "RetryEscalationWorkflow":     "Retry escalation",
        "ScheduledMaintenanceWorkflow":"Maintenance",
    }

    # Also check for child workflows
    child_types = {
        "PDFConversionWorkflow", "DOCXConversionWorkflow",
        "XLSXConversionWorkflow", "PPTXConversionWorkflow",
        "HTMLConversionWorkflow", "TXTConversionWorkflow",
        "ImageConversionWorkflow",
    }

    subhead("Checking expected workflow types:")
    for wf_type, desc in expected_types.items():
        if wf_type in by_type:
            count = len(by_type[wf_type])
            completed = sum(1 for w in by_type[wf_type] if w.get("status")=="COMPLETED")
            ok(f"{wf_type}  ({completed}/{count} completed)  – {desc}")
        else:
            fail(f"{wf_type} NOT found in Temporal  – {desc}")

    subhead("Child workflows (per-document-type):")
    found_children = set(by_type.keys()) & child_types
    if found_children:
        for ct in sorted(found_children):
            count = len(by_type[ct])
            ok(f"{ct}  ({count} executions)")
    else:
        fail("No child workflows found – conversion may have used direct pipeline")

    # Overall count
    total = len(workflows)
    completed = sum(1 for w in workflows if w.get("status")=="COMPLETED")
    failed = sum(1 for w in workflows if w.get("status")=="FAILED")
    running = sum(1 for w in workflows if w.get("status")=="RUNNING")

    subhead(f"Dashboard totals: {total} workflows  |  "
            f"{G}{completed} completed{N}  |  {R}{failed} failed{N}  |  {Y}{running} running{N}")

    if completed > 0:
        ok(f"Temporal dashboard has {completed} completed workflows visible")
    if failed > 0:
        fail(f"{failed} workflows failed")


def test_s3_outputs():
    """Verify all output files exist in S3."""
    header("12. S3 Output Verification")

    keys = s3_ls(OUT_BUCKET, "")
    print(f"  Total objects in s3://{OUT_BUCKET}/: {len(keys)}\n")

    prefixes = {}
    for k in keys:
        prefix = k.split("/")[0]
        prefixes[prefix] = prefixes.get(prefix, 0) + 1

    for prefix, count in sorted(prefixes.items()):
        print(f"    {prefix}/  →  {count} files")

    if len(keys) > 0:
        ok(f"S3 output bucket has {len(keys)} objects")
    else:
        fail("S3 output bucket is empty")

    # Sample one converted text
    converted = [k for k in keys if k.startswith("converted/") and k.endswith(".txt")]
    if converted:
        text = s3_get(OUT_BUCKET, converted[0])
        if text:
            ok(f"Sample output readable: {converted[0]}")
            preview = text[:200].replace("\n"," ↵ ")
            print(f"    {D}{preview}...{N}")


# ═════════════════════════════════════════════════════════════════════════════
# MAIN
# ═════════════════════════════════════════════════════════════════════════════

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--api-base", default=API)
    parser.add_argument("--s3-endpoint", default=S3_EP)
    args = parser.parse_args()

    global API, S3_EP
    API = args.api_base
    S3_EP = args.s3_endpoint
    AWS["endpoint_url"] = S3_EP

    header("Document Conversion Service – Full Verification")
    print(f"  API:         {API}")
    print(f"  S3:          {S3_EP}")
    print(f"  Temporal UI: http://localhost:8088")
    print(f"  Time:        {time.strftime('%Y-%m-%d %H:%M:%S')}")

    # Run all tests in order
    if not test_infrastructure():
        print(f"\n  {R}Infrastructure not ready – aborting{N}\n")
        return 1

    tests = [
        test_upload_samples,
        test_direct_upload,
        test_s3_conversion,
        test_pipeline_workflow,
        test_s3_watch_workflow,
        test_multi_format_workflow,
        test_webhook_workflow,
        test_retry_escalation,
        test_maintenance,
        test_temporal_dashboard,
        test_s3_outputs,
    ]

    for test_fn in tests:
        try:
            test_fn()
        except KeyboardInterrupt:
            print("\n  Interrupted."); break
        except Exception as e:
            fail(f"Test {test_fn.__name__} crashed: {e}")
            traceback.print_exc()

    # ── Summary ──────────────────────────────────────────────────────────
    header("VERIFICATION SUMMARY")

    print(f"  {'Test':<55s}  Status")
    print(f"  {'─'*55}  {'─'*8}")
    for status, name in RESULTS:
        color = G if status=="PASS" else R
        print(f"  {name:<55s}  {color}{status}{N}")

    print(f"\n  {B}Total: {PASS_COUNT} passed, {FAIL_COUNT} failed{N}\n")

    if FAIL_COUNT == 0:
        print(f"  {G}{B}{'═'*50}")
        print(f"  ★  ALL VERIFICATIONS PASSED  ★")
        print(f"  {'═'*50}{N}")
        print(f"\n  Open Temporal UI to see all workflows: {B}http://localhost:8088{N}\n")
    else:
        print(f"  {R}{B}{'═'*50}")
        print(f"  ✗  {FAIL_COUNT} VERIFICATION(S) FAILED")
        print(f"  {'═'*50}{N}")

    return 0 if FAIL_COUNT == 0 else 1


if __name__ == "__main__":
    sys.exit(main())
