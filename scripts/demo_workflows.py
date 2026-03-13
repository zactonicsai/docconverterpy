#!/usr/bin/env python3
"""
demo_workflows.py – End-to-end demo of all document conversion workflows.

This script:
  1. Creates sample documents (PDF, DOCX, XLSX, PPTX, TXT, HTML, image)
  2. Uploads them to S3 (LocalStack) in the input bucket
  3. Triggers every workflow type via the REST API
  4. Polls workflow status until completion
  5. Downloads and displays the converted text from S3
  6. Shows a summary of all results

Prerequisites:
  - Docker stack running:  docker compose up -d
  - Wait ~60s for all services to be healthy
  - pip install boto3 requests reportlab python-docx openpyxl python-pptx Pillow

Usage:
  python scripts/demo_workflows.py                    # run all demos
  python scripts/demo_workflows.py --only pipeline    # run one workflow
  python scripts/demo_workflows.py --api-base http://myhost:8080
"""

import argparse
import io
import json
import os
import sys
import textwrap
import time

import boto3
import requests

# ─── Configuration ───────────────────────────────────────────────────────────

API_BASE = "http://localhost:8080"
S3_ENDPOINT = "http://localhost:4566"
INPUT_BUCKET = "docconv-input"
OUTPUT_BUCKET = "docconv-output"
AWS_KWARGS = {
    "endpoint_url": S3_ENDPOINT,
    "aws_access_key_id": "test",
    "aws_secret_access_key": "test",
    "region_name": "us-east-1",
}

# ─── Colors ──────────────────────────────────────────────────────────────────

G = "\033[92m"; R = "\033[91m"; C = "\033[96m"; Y = "\033[93m"
B = "\033[1m"; D = "\033[2m"; N = "\033[0m"

def ok(msg):   print(f"  {G}✓{N} {msg}")
def fail(msg): print(f"  {R}✗{N} {msg}")
def info(msg): print(f"  {C}ℹ{N} {msg}")
def warn(msg): print(f"  {Y}⚠{N} {msg}")
def header(title):
    print(f"\n{B}{C}{'═'*70}")
    print(f"  {title}")
    print(f"{'═'*70}{N}\n")
def subhead(title):
    print(f"\n  {B}{title}{N}")

# ─── S3 Helpers ──────────────────────────────────────────────────────────────

def s3_client():
    return boto3.client("s3", **AWS_KWARGS)

def s3_upload_bytes(bucket, key, data, content_type="application/octet-stream"):
    s3_client().put_object(Bucket=bucket, Key=key, Body=data, ContentType=content_type)
    ok(f"Uploaded s3://{bucket}/{key}  ({len(data):,} bytes)")

def s3_upload_file(bucket, key, filepath):
    s3_client().upload_file(filepath, bucket, key)
    size = os.path.getsize(filepath)
    ok(f"Uploaded s3://{bucket}/{key}  ({size:,} bytes)")

def s3_download_text(bucket, key):
    try:
        resp = s3_client().get_object(Bucket=bucket, Key=key)
        return resp["Body"].read().decode("utf-8")
    except Exception as e:
        return f"[download failed: {e}]"

def s3_list_keys(bucket, prefix):
    try:
        resp = s3_client().list_objects_v2(Bucket=bucket, Prefix=prefix)
        return [obj["Key"] for obj in resp.get("Contents", [])]
    except Exception:
        return []

def ensure_buckets():
    """Create input/output buckets if they don't exist."""
    client = s3_client()
    for bucket in [INPUT_BUCKET, OUTPUT_BUCKET]:
        try:
            client.head_bucket(Bucket=bucket)
        except Exception:
            client.create_bucket(Bucket=bucket)
            ok(f"Created bucket: {bucket}")

# ─── Sample File Generators ─────────────────────────────────────────────────

def make_sample_pdf():
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas as rl
    buf = io.BytesIO()
    c = rl.Canvas(buf, pagesize=letter)
    c.setFont("Helvetica", 14)
    c.drawString(72, 700, "Document Conversion Service - Sample PDF")
    c.setFont("Helvetica", 11)
    y = 660
    for line in textwrap.wrap(
        "This PDF document contains two pages of text content. "
        "The conversion service extracts text from each page using pdfplumber. "
        "For scanned documents, it falls back to OCR using Tesseract. "
        "Embedded images are also extracted and OCR'd separately.", width=80
    ):
        c.drawString(72, y, line); y -= 16
    c.drawString(72, y - 20, "Page 1 of 2")
    c.showPage()
    c.setFont("Helvetica", 11)
    c.drawString(72, 700, "This is page two with additional content.")
    c.drawString(72, 680, "Tables, charts, and images would be extracted here.")
    c.drawString(72, 660, "Page 2 of 2")
    c.showPage()
    c.save()
    return buf.getvalue()

def make_sample_docx():
    from docx import Document
    doc = Document()
    doc.add_heading("Quarterly Report", level=1)
    doc.add_paragraph(
        "This Word document demonstrates DOCX conversion. "
        "The service extracts paragraphs and tables using python-docx."
    )
    table = doc.add_table(rows=3, cols=3)
    for i, header in enumerate(["Product", "Q1 Revenue", "Q2 Revenue"]):
        table.cell(0, i).text = header
    for i, row in enumerate([["Widget A", "$1,500", "$2,300"], ["Widget B", "$800", "$1,200"]], 1):
        for j, val in enumerate(row):
            table.cell(i, j).text = val
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()

def make_sample_xlsx():
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Sales Data"
    ws.append(["Month", "Revenue", "Expenses", "Profit"])
    for month, rev, exp in [("Jan", 50000, 35000), ("Feb", 62000, 38000),
                             ("Mar", 58000, 36000), ("Apr", 71000, 42000)]:
        ws.append([month, rev, exp, rev - exp])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()

def make_sample_pptx():
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Project Status Update"
    slide.placeholders[1].text = (
        "Key accomplishments this quarter:\n"
        "• Launched document conversion service\n"
        "• Integrated Temporal.io workflows\n"
        "• Added 10 format converters"
    )
    notes = slide.notes_slide
    notes.notes_text_frame.text = "Speaker note: mention the Temporal integration."
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Next Steps"
    slide2.placeholders[1].text = "Continue scaling the pipeline and adding connectors."
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

def make_sample_html():
    return b"""<!DOCTYPE html>
<html><head><title>Sample Report</title></head>
<body>
<h1>Annual Technology Review</h1>
<p>This HTML document demonstrates web page conversion. The service strips
all HTML tags, scripts, and styles, returning only the visible text content.</p>
<h2>Key Findings</h2>
<ul>
<li>Cloud adoption increased by 34% year over year</li>
<li>Container workloads now represent 60% of production deployments</li>
<li>Temporal.io adoption growing rapidly for durable execution</li>
</ul>
<table><tr><th>Metric</th><th>2024</th><th>2025</th></tr>
<tr><td>Uptime</td><td>99.9%</td><td>99.99%</td></tr>
<tr><td>Latency</td><td>120ms</td><td>45ms</td></tr></table>
<script>var x = 1;</script>
</body></html>"""

def make_sample_txt():
    return (
        "Document Conversion Service - Plain Text Sample\n"
        "================================================\n\n"
        "This is a plain text file used to test the TXT converter.\n"
        "The converter detects encoding (UTF-8, Latin-1, etc.) and reads\n"
        "the file in configurable chunks for memory efficiency.\n\n"
        "Line counts, word counts, and character statistics can be\n"
        "computed by the text analytics activity.\n"
    ).encode("utf-8")

def make_sample_image():
    from PIL import Image, ImageDraw, ImageFont
    img = Image.new("RGB", (600, 250), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 24)
        small = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 18)
    except (IOError, OSError):
        font = ImageFont.load_default()
        small = font
    draw.text((30, 20), "Invoice #2025-0042", fill=(0, 0, 0), font=font)
    draw.text((30, 60), "Customer: Acme Corporation", fill=(0, 0, 0), font=small)
    draw.text((30, 90), "Amount: $12,450.00", fill=(0, 0, 0), font=small)
    draw.text((30, 120), "Date: March 13, 2025", fill=(0, 0, 0), font=small)
    draw.text((30, 160), "OCR extracts text from images", fill=(80, 80, 80), font=small)
    draw.text((30, 190), "using Tesseract engine", fill=(80, 80, 80), font=small)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()

# ─── API Helpers ─────────────────────────────────────────────────────────────

def api_post(endpoint, json_body=None, files=None, data=None):
    url = f"{API_BASE}{endpoint}"
    try:
        if files:
            resp = requests.post(url, files=files, data=data, timeout=120)
        else:
            resp = requests.post(url, json=json_body, timeout=120)
        return resp.status_code, resp.json()
    except Exception as e:
        return 0, {"error": str(e)}

def api_get(endpoint):
    try:
        resp = requests.get(f"{API_BASE}{endpoint}", timeout=30)
        return resp.status_code, resp.json()
    except Exception as e:
        return 0, {"error": str(e)}

def poll_workflow(workflow_id, max_wait=120):
    """Poll a workflow until it completes or times out."""
    start = time.time()
    last_status = ""
    while time.time() - start < max_wait:
        code, data = api_get(f"/workflow/{workflow_id}/status")
        if code == 200:
            status = data.get("status", "")
            custom = data.get("custom_status", "")
            step = data.get("current_step", "")
            display = custom or status

            if display != last_status:
                print(f"    {D}[{time.time()-start:.0f}s]{N} {display}  step={step}")
                last_status = display

            if status in ("COMPLETED", "FAILED", "TERMINATED", "CANCELLED", "TIMED_OUT"):
                return status, data
        elif code == 404:
            # Workflow may not be registered yet
            pass

        time.sleep(2)

    return "TIMEOUT", {}

def show_result(label, code, data):
    success = data.get("success", code == 200)
    chars = data.get("characters_extracted", data.get("total_chars", ""))
    bucket = data.get("output_bucket", data.get("text_bucket", ""))
    key = data.get("output_key", data.get("text_key", ""))

    if success:
        parts = [f"{G}SUCCESS{N}"]
        if chars: parts.append(f"chars={chars}")
        if key:   parts.append(f"→ s3://{bucket}/{key}")
        ok(f"{label}:  {' | '.join(parts)}")
    else:
        err = data.get("error", "unknown")
        fail(f"{label}:  {R}FAILED{N} – {err}")

def show_s3_text(bucket, key, max_chars=500):
    """Download and display converted text from S3."""
    text = s3_download_text(bucket, key)
    if text.startswith("[download"):
        warn(f"  Could not download: {text}")
        return
    print(f"    {D}── s3://{bucket}/{key} ──{N}")
    lines = text[:max_chars].splitlines()
    for line in lines[:15]:
        print(f"    {line}")
    if len(text) > max_chars:
        print(f"    {Y}... ({len(text) - max_chars:,} more chars){N}")
    print()

# ═════════════════════════════════════════════════════════════════════════════
# WORKFLOW DEMOS
# ═════════════════════════════════════════════════════════════════════════════

def demo_health():
    header("0. Health Check")
    code, data = api_get("/health")
    if code == 200:
        ok(f"API is running")
        for k, v in data.items():
            print(f"    {k}: {v}")
    else:
        fail(f"API not reachable at {API_BASE}")
        sys.exit(1)

def demo_upload_samples():
    header("1. Upload Sample Files to S3")
    ensure_buckets()

    samples = {
        "docs/report.pdf":      ("application/pdf",     make_sample_pdf()),
        "docs/report.docx":     ("application/vnd.openxmlformats", make_sample_docx()),
        "docs/data.xlsx":       ("application/vnd.openxmlformats", make_sample_xlsx()),
        "docs/slides.pptx":     ("application/vnd.openxmlformats", make_sample_pptx()),
        "docs/page.html":       ("text/html",            make_sample_html()),
        "docs/notes.txt":       ("text/plain",           make_sample_txt()),
        "docs/invoice.png":     ("image/png",            make_sample_image()),
        # Files for inbox/ folder watch
        "inbox/report.pdf":     ("application/pdf",      make_sample_pdf()),
        "inbox/data.xlsx":      ("application/vnd.openxmlformats", make_sample_xlsx()),
        "inbox/notes.txt":      ("text/plain",           make_sample_txt()),
        "inbox/invoice.png":    ("image/png",            make_sample_image()),
    }

    for key, (ctype, data) in samples.items():
        s3_upload_bytes(INPUT_BUCKET, key, data, ctype)

    info(f"Total: {len(samples)} files uploaded to s3://{INPUT_BUCKET}/")

def demo_api_upload():
    header("2. REST API – Direct File Upload")
    info("Upload files directly via POST /convert/upload (bypasses Temporal)")

    uploads = [
        ("report.pdf",   "pdf",   make_sample_pdf()),
        ("report.docx",  "docx",  make_sample_docx()),
        ("data.xlsx",    "xlsx",  make_sample_xlsx()),
        ("slides.pptx",  "pptx",  make_sample_pptx()),
        ("page.html",    "html",  make_sample_html()),
        ("notes.txt",    "txt",   make_sample_txt()),
        ("invoice.png",  "image", make_sample_image()),
    ]

    for filename, doc_type, data in uploads:
        code, result = api_post(
            "/convert/upload",
            files={"file": (filename, data)},
            data={"document_type": doc_type},
        )
        show_result(f"Upload {filename}", code, result)
        if result.get("success") and result.get("output_key"):
            show_s3_text(result["output_bucket"], result["output_key"])

def demo_basic_conversion():
    header("3. Basic Conversion Workflow (URL → text)")
    info("POST /convert/job with location_type=url → Temporal workflow")

    code, result = api_post("/convert/job", {
        "job_id": "demo-url-001",
        "document_type": "html",
        "location_type": "url",
        "url": "https://example.com",
    })
    show_result("URL conversion", code, result)
    if result.get("output_key"):
        show_s3_text(result["output_bucket"], result["output_key"])

def demo_s3_conversion():
    header("4. S3 Conversion Workflow (S3 → text)")
    info("POST /convert/job with location_type=s3 → Temporal workflow")

    for key, doc_type, label in [
        ("docs/report.pdf",  "pdf",   "PDF from S3"),
        ("docs/report.docx", "docx",  "DOCX from S3"),
        ("docs/data.xlsx",   "xlsx",  "XLSX from S3"),
        ("docs/slides.pptx", "pptx",  "PPTX from S3"),
        ("docs/page.html",   "html",  "HTML from S3"),
        ("docs/notes.txt",   "txt",   "TXT from S3"),
        ("docs/invoice.png", "image", "Image from S3"),
    ]:
        job_id = f"demo-s3-{doc_type}-001"
        code, result = api_post("/convert/job", {
            "job_id": job_id,
            "document_type": doc_type,
            "location_type": "s3",
            "s3_bucket": INPUT_BUCKET,
            "s3_key": key,
            "s3_endpoint_url": "http://localstack:4566",
        })
        show_result(label, code, result)
        if result.get("output_key"):
            show_s3_text(result["output_bucket"], result["output_key"])

def demo_pipeline():
    header("5. Pipeline Workflow (validate → convert → analyze → metadata)")
    info("POST /workflow/pipeline → DocumentPipelineWorkflow")

    code, result = api_post("/workflow/pipeline", {
        "job_id": "demo-pipe-001",
        "document_type": "pdf",
        "location_type": "s3",
        "s3_bucket": INPUT_BUCKET,
        "s3_key": "docs/report.pdf",
        "s3_endpoint_url": "http://localstack:4566",
        "output_prefix": "pipeline/",
        "enable_validation": True,
        "enable_analytics": True,
        "enable_metadata_output": True,
    })

    if result.get("workflow_id"):
        ok(f"Workflow started: {result['workflow_id']}")
        subhead("Polling workflow status...")
        status, _ = poll_workflow(result["workflow_id"])
        ok(f"Final status: {status}")

        # Show outputs
        subhead("Pipeline outputs:")
        for key in s3_list_keys(OUTPUT_BUCKET, "pipeline/demo-pipe-001/"):
            info(f"s3://{OUTPUT_BUCKET}/{key}")
            if key.endswith(".txt"):
                show_s3_text(OUTPUT_BUCKET, key)
            elif key.endswith(".json"):
                meta = s3_download_text(OUTPUT_BUCKET, key)
                print(f"    {meta[:600]}")
                print()
    else:
        fail(f"Pipeline start failed: {result}")

def demo_s3_watch():
    header("6. S3 Folder Watch Workflow (scan inbox/ → convert all)")
    info("POST /workflow/s3-watch → S3FolderWatchWorkflow")

    code, result = api_post("/workflow/s3-watch", {
        "bucket": INPUT_BUCKET,
        "prefix": "inbox/",
        "output_prefix": "watch-output/",
        "s3_endpoint_url": "http://localstack:4566",
        "move_processed_to": "processed/",
        "max_files": 50,
    })

    if result.get("workflow_id"):
        ok(f"Workflow started: {result['workflow_id']}")
        subhead("Polling workflow status...")
        status, _ = poll_workflow(result["workflow_id"], max_wait=180)
        ok(f"Final status: {status}")

        subhead("Watch outputs:")
        for key in s3_list_keys(OUTPUT_BUCKET, "watch-output/"):
            info(f"s3://{OUTPUT_BUCKET}/{key}")

        subhead("Processed (moved from inbox/):")
        for key in s3_list_keys(INPUT_BUCKET, "processed/"):
            info(f"s3://{INPUT_BUCKET}/{key}")
    else:
        fail(f"S3 watch start failed: {result}")

def demo_multi_format():
    header("7. Multi-Format Output (full text + pages + metadata + analytics)")
    info("POST /workflow/multi-format → MultiFormatOutputWorkflow")

    code, result = api_post("/workflow/multi-format", {
        "job_id": "demo-multi-001",
        "document_type": "pdf",
        "location_type": "s3",
        "s3_bucket": INPUT_BUCKET,
        "s3_key": "docs/report.pdf",
        "s3_endpoint_url": "http://localstack:4566",
        "output_prefix": "multi/",
        "produce_full_text": True,
        "produce_pages": True,
        "produce_metadata_json": True,
        "produce_analytics": True,
    })

    if result.get("workflow_id"):
        ok(f"Workflow started: {result['workflow_id']}")
        subhead("Polling workflow status...")
        status, _ = poll_workflow(result["workflow_id"])
        ok(f"Final status: {status}")

        subhead("Multi-format outputs:")
        for key in sorted(s3_list_keys(OUTPUT_BUCKET, "multi/demo-multi-001/")):
            info(f"s3://{OUTPUT_BUCKET}/{key}")
        # Show full text
        show_s3_text(OUTPUT_BUCKET, "multi/demo-multi-001/full_text.txt")
        # Show metadata
        meta = s3_download_text(OUTPUT_BUCKET, "multi/demo-multi-001/metadata.json")
        if not meta.startswith("[download"):
            subhead("Metadata JSON:")
            print(f"    {meta[:600]}")
            print()
    else:
        fail(f"Multi-format start failed: {result}")

def demo_webhook():
    header("8. Webhook Notification Workflow")
    info("POST /workflow/webhook-convert → WebhookNotificationWorkflow")
    info("Using https://httpbin.org/post as webhook receiver")

    code, result = api_post("/workflow/webhook-convert", {
        "job_id": "demo-wh-001",
        "document_type": "html",
        "location_type": "url",
        "url": "https://example.com",
        "on_start_webhook": "https://httpbin.org/post",
        "on_complete_webhook": "https://httpbin.org/post",
        "on_failure_webhook": "https://httpbin.org/post",
    })

    if result.get("workflow_id"):
        ok(f"Workflow started: {result['workflow_id']}")
        subhead("Polling workflow status...")
        status, _ = poll_workflow(result["workflow_id"])
        ok(f"Final status: {status}")
    else:
        fail(f"Webhook workflow start failed: {result}")

def demo_retry_escalation():
    header("9. Retry Escalation Workflow (standard → enhanced OCR)")
    info("POST /workflow/retry-escalation → RetryEscalationWorkflow")

    code, result = api_post("/workflow/retry-escalation", {
        "job_id": "demo-esc-001",
        "document_type": "pdf",
        "location_type": "s3",
        "s3_bucket": INPUT_BUCKET,
        "s3_key": "docs/report.pdf",
        "s3_endpoint_url": "http://localstack:4566",
        "min_chars_threshold": 50,
        "escalation_dpi": 400,
    })

    if result.get("workflow_id"):
        ok(f"Workflow started: {result['workflow_id']}")
        subhead("Polling workflow status...")
        status, _ = poll_workflow(result["workflow_id"], max_wait=180)
        ok(f"Final status: {status}")
    else:
        fail(f"Escalation start failed: {result}")

def demo_maintenance():
    header("10. Scheduled Maintenance Workflow")
    info("POST /workflow/maintenance → ScheduledMaintenanceWorkflow")

    code, result = api_post("/workflow/maintenance", {
        "tmp_dir": "/tmp/docconv",
        "max_age_hours": 24,
    })

    if result.get("workflow_id"):
        ok(f"Workflow started: {result['workflow_id']}")
        subhead("Polling workflow status...")
        status, _ = poll_workflow(result["workflow_id"], max_wait=30)
        ok(f"Final status: {status}")
    else:
        fail(f"Maintenance start failed: {result}")

def demo_workflow_queries():
    header("11. Query Workflow Status & List Recent")

    subhead("Recent workflows:")
    code, data = api_get("/workflows/recent?limit=10")
    if code == 200:
        for wf in data.get("workflows", []):
            status_color = G if wf["status"] == "COMPLETED" else R if wf["status"] == "FAILED" else Y
            print(f"    {status_color}{wf['status']:12s}{N}  {wf['id']}")
    else:
        warn(f"Could not list workflows: {data}")

def demo_show_all_outputs():
    header("12. All Converted Outputs in S3")

    subhead(f"s3://{OUTPUT_BUCKET}/")
    keys = s3_list_keys(OUTPUT_BUCKET, "")
    if keys:
        for key in sorted(keys):
            size = ""
            try:
                resp = s3_client().head_object(Bucket=OUTPUT_BUCKET, Key=key)
                size = f"  ({resp['ContentLength']:,} bytes)"
            except Exception:
                pass
            print(f"    {key}{D}{size}{N}")
    else:
        warn("No output files found")

# ═════════════════════════════════════════════════════════════════════════════
# MAIN
# ═════════════════════════════════════════════════════════════════════════════

DEMOS = {
    "health":     demo_health,
    "upload":     demo_upload_samples,
    "api":        demo_api_upload,
    "basic":      demo_basic_conversion,
    "s3":         demo_s3_conversion,
    "pipeline":   demo_pipeline,
    "s3watch":    demo_s3_watch,
    "multi":      demo_multi_format,
    "webhook":    demo_webhook,
    "escalation": demo_retry_escalation,
    "maintenance":demo_maintenance,
    "queries":    demo_workflow_queries,
    "outputs":    demo_show_all_outputs,
}

def main():
    parser = argparse.ArgumentParser(description="DocConv Workflow Demo")
    parser.add_argument("--api-base", default=API_BASE, help="API base URL")
    parser.add_argument("--s3-endpoint", default=S3_ENDPOINT, help="S3 endpoint URL")
    parser.add_argument("--only", help=f"Run only one demo: {', '.join(DEMOS.keys())}")
    parser.add_argument("--list", action="store_true", help="List available demos")
    args = parser.parse_args()

    global API_BASE, S3_ENDPOINT
    API_BASE = args.api_base
    S3_ENDPOINT = args.s3_endpoint
    AWS_KWARGS["endpoint_url"] = S3_ENDPOINT

    if args.list:
        print("Available demos:")
        for name in DEMOS:
            print(f"  {name}")
        return

    header("Document Conversion Service – Workflow Demo")
    print(f"  API:  {API_BASE}")
    print(f"  S3:   {S3_ENDPOINT}")
    print(f"  View Temporal UI: http://localhost:8088")

    if args.only:
        if args.only not in DEMOS:
            fail(f"Unknown demo: {args.only}. Use --list to see options.")
            return
        demo_health()
        if args.only != "health":
            demo_upload_samples()
            DEMOS[args.only]()
    else:
        for name, fn in DEMOS.items():
            try:
                fn()
            except KeyboardInterrupt:
                print("\n  Interrupted.")
                break
            except Exception as e:
                fail(f"Demo '{name}' crashed: {e}")

    header("Demo Complete!")
    print(f"  Temporal UI:       http://localhost:8088")
    print(f"  API Test Console:  {API_BASE}")
    print(f"  RabbitMQ UI:       http://localhost:15672  (docconv/docconv)")
    print()

if __name__ == "__main__":
    main()
