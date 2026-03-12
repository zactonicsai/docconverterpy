"""
PPTX → Text converter.

Extracts text from each slide shape and speaker notes.
Yields text one slide at a time.
"""

import logging
from typing import Generator

from pptx import Presentation

logger = logging.getLogger(__name__)


def convert_to_text(file_path: str) -> Generator[str, None, None]:
    logger.info("PPTX convert: %s", file_path)
    prs = Presentation(file_path)

    for slide_num, slide in enumerate(prs.slides, start=1):
        parts: list[str] = [f"--- Slide {slide_num} ---"]

        # Shape text
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)

            # Table shapes
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    parts.append(" | ".join(cells))

        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[NOTES] {notes}")

        if len(parts) > 1:  # more than just the header
            yield "\n".join(parts) + "\n"
